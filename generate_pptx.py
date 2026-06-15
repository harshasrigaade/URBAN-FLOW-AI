import subprocess
import sys
import os

# Step 1: Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()

# Set slide width and height for widescreen 16:9
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Brand Colors
DEEP_NAVY = RGBColor(11, 19, 43)      # #0B132B
DARK_CARD = RGBColor(28, 37, 65)      # #1C2541
BORDER_BLUE = RGBColor(58, 80, 107)   # #3A506B
NEON_CYAN = RGBColor(111, 255, 233)   # #6FFFE9
SOFT_TEAL = RGBColor(91, 192, 190)    # #5BC0BE
EMERALD_GREEN = RGBColor(16, 185, 129) # #10B981
CRIMSON_RED = RGBColor(239, 68, 68)   # #EF4444
WARNING_GOLD = RGBColor(245, 158, 11) # #F59E0B
WHITE = RGBColor(248, 250, 252)       # #F8FAFC
SLATE_GRAY = RGBColor(148, 163, 184)  # #94A3B8

def apply_background(slide):
    """Fills background with Deep Navy color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DEEP_NAVY

def add_header(slide, title_text, category_text=None):
    """Adds a standard header with category tag and title"""
    # Category tag (optional)
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Century Gothic'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = SOFT_TEAL
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Century Gothic'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    
    # Decorative line under header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_BLUE
    line.line.fill.background() # No border

def create_card(slide, left, top, width, height, fill_color, border_color=None):
    """Creates a stylized background card with optional border color"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

# ----------------- SLIDE 1: COVER SLIDE -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)

# Subtitle Category
cat_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.4))
tf_cat = cat_box.text_frame
p_cat = tf_cat.paragraphs[0]
p_cat.text = "AI-POWERED SMART MOBILITY ECOSYSTEM"
p_cat.font.name = 'Century Gothic'
p_cat.font.size = Pt(14)
p_cat.font.bold = True
p_cat.font.color.rgb = NEON_CYAN
p_cat.alignment = PP_ALIGN.CENTER

# Main Title
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2))
tf_title = title_box.text_frame
p_title = tf_title.paragraphs[0]
p_title.text = "URBANFLOW AI"
p_title.font.name = 'Century Gothic'
p_title.font.size = Pt(64)
p_title.font.bold = True
p_title.font.color.rgb = WHITE
p_title.alignment = PP_ALIGN.CENTER

# Tagline
tag_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.6))
tf_tag = tag_box.text_frame
p_tag = tf_tag.paragraphs[0]
p_tag.text = "Travel Faster, Safer, Cheaper, and Greener."
p_tag.font.name = 'Century Gothic'
p_tag.font.size = Pt(18)
p_tag.font.color.rgb = SLATE_GRAY
p_tag.alignment = PP_ALIGN.CENTER

# Team info box card
create_card(slide, Inches(3.66), Inches(4.5), Inches(6.0), Inches(1.8), DARK_CARD, BORDER_BLUE)

team_box = slide.shapes.add_textbox(Inches(3.8), Inches(4.6), Inches(5.72), Inches(1.6))
tf_team = team_box.text_frame
tf_team.word_wrap = True

p_team_header = tf_team.paragraphs[0]
p_team_header.text = "DEVELOPED FOR HACKATHON BY:"
p_team_header.font.name = 'Century Gothic'
p_team_header.font.size = Pt(10)
p_team_header.font.bold = True
p_team_header.font.color.rgb = SOFT_TEAL
p_team_header.alignment = PP_ALIGN.CENTER

p_mem1 = tf_team.add_paragraph()
p_mem1.text = "G. Harsha Sri  |  AI Specialist"
p_mem1.font.name = 'Century Gothic'
p_mem1.font.size = Pt(14)
p_mem1.font.bold = True
p_mem1.font.color.rgb = WHITE
p_mem1.alignment = PP_ALIGN.CENTER
p_mem1.space_before = Pt(10)

p_mem2 = tf_team.add_paragraph()
p_mem2.text = "P. Sindhu Reddy  |  Full-Stack Architect"
p_mem2.font.name = 'Century Gothic'
p_mem2.font.size = Pt(14)
p_mem2.font.bold = True
p_mem2.font.color.rgb = WHITE
p_mem2.alignment = PP_ALIGN.CENTER
p_mem2.space_before = Pt(5)


# ----------------- SLIDE 2: THE PROBLEM STATEMENT -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "The Fragmented & Unsafe Urban Journey", "The Problem Statement")

# Add a description text box
desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
tf_desc = desc_box.text_frame
p_desc = tf_desc.paragraphs[0]
p_desc.text = "Modern commuters deal with a high-friction transit cycle. Information is scattered, safety risks are elevated, and environmental impact is ignored."
p_desc.font.name = 'Century Gothic'
p_desc.font.size = Pt(14)
p_desc.font.color.rgb = SLATE_GRAY

# 4 column cards for 4 core problems
col_width = Inches(2.7)
col_height = Inches(4.0)
gap = Inches(0.3)
left_margin = Inches(0.8)
top_margin = Inches(2.1)

problems = [
    {
        "title": "SEVERE CONGESTION",
        "label": "Time Drainage",
        "desc": "Commuters waste hundreds of hours in gridlock daily. Traditional maps only offer reactive notifications after gridlocks occur.",
        "color": CRIMSON_RED
    },
    {
        "title": "SAFETY BLINDSPOTS",
        "label": "Physical Vulnerability",
        "desc": "Poorly lit sectors, high-crime corridors, and accident hotspots are not integrated into route recommendations by generic mapping tools.",
        "color": WARNING_GOLD
    },
    {
        "title": "RISING COMMUTE COSTS",
        "label": "Economic Burden",
        "desc": "commuters struggle to compare fares across cabs, buses, metro lines, and auto-rickshaws to optimize their daily travel budget.",
        "color": CRIMSON_RED
    },
    {
        "title": "CARBON FOOTPRINT",
        "label": "Ecological Stress",
        "desc": "There is no simplified tool to track carbon offsets or reward green travel alternatives, keeping public transit adoption low.",
        "color": WARNING_GOLD
    }
]

for idx, prob in enumerate(problems):
    col_left = left_margin + idx * (col_width + gap)
    create_card(slide, col_left, top_margin, col_width, col_height, DARK_CARD, prob["color"])
    
    # Text inside card
    box = slide.shapes.add_textbox(col_left + Inches(0.15), top_margin + Inches(0.15), col_width - Inches(0.3), col_height - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    
    p_header = tf.paragraphs[0]
    p_header.text = prob["title"]
    p_header.font.name = 'Century Gothic'
    p_header.font.size = Pt(14)
    p_header.font.bold = True
    p_header.font.color.rgb = WHITE
    
    p_label = tf.add_paragraph()
    p_label.text = prob["label"].upper()
    p_label.font.name = 'Century Gothic'
    p_label.font.size = Pt(9)
    p_label.font.bold = True
    p_label.font.color.rgb = prob["color"]
    p_label.space_after = Pt(14)
    
    p_desc = tf.add_paragraph()
    p_desc.text = prob["desc"]
    p_desc.font.name = 'Century Gothic'
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = SLATE_GRAY


# ----------------- SLIDE 3: THE SOLUTION -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "UrbanFlow AI: The Smart Mobility Ecosystem", "The Unified Platform")

# Left text description
left_col = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.2), Inches(4.8))
tf_left = left_col.text_frame
tf_left.word_wrap = True

p_sol_h = tf_left.paragraphs[0]
p_sol_h.text = "A Single digital corridor solving multi-objective travel routing"
p_sol_h.font.name = 'Century Gothic'
p_sol_h.font.size = Pt(22)
p_sol_h.font.bold = True
p_sol_h.font.color.rgb = WHITE
p_sol_h.space_after = Pt(14)

p_sol_d = tf_left.add_paragraph()
p_sol_d.text = "UrbanFlow AI acts as a smart transit assistant, addressing multiple travel criteria simultaneously in a single application."
p_sol_d.font.name = 'Century Gothic'
p_sol_d.font.size = Pt(13)
p_sol_d.font.color.rgb = SLATE_GRAY
p_sol_d.space_after = Pt(10)

p_sol_d2 = tf_left.add_paragraph()
p_sol_d2.text = "Rather than switching between safety apps, taxi bookings, and carbon trackers, commuters access an integrated dynamic client dashboard powered by live geospatial analysis."
p_sol_d2.font.name = 'Century Gothic'
p_sol_d2.font.size = Pt(13)
p_sol_d2.font.color.rgb = SLATE_GRAY

# Right diagram card
create_card(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.8), DARK_CARD, BORDER_BLUE)

diag_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.9), Inches(5.6), Inches(4.4))
tf_diag = diag_box.text_frame
tf_diag.word_wrap = True

p_diag_title = tf_diag.paragraphs[0]
p_diag_title.text = "THE INTEGRATED PLATFORM EQUATION"
p_diag_title.font.name = 'Century Gothic'
p_diag_title.font.size = Pt(11)
p_diag_title.font.bold = True
p_diag_title.font.color.rgb = NEON_CYAN
p_diag_title.space_after = Pt(14)

modules = [
    ("MAPS & ROUTING ENGINE", "OpenStreetMap / OSRM Graph Router", SOFT_TEAL),
    ("SAFETY AI + EMERGENCY SOS", "Guardian Safety Intelligence Suite", CRIMSON_RED),
    ("CARBON ACCOUNTING & XP REWARDS", "Sustained Commuting Gamification System", EMERALD_GREEN)
]

for title, desc, color in modules:
    p_m_title = tf_diag.add_paragraph()
    p_m_title.text = f"▪ {title}"
    p_m_title.font.name = 'Century Gothic'
    p_m_title.font.size = Pt(12)
    p_m_title.font.bold = True
    p_m_title.font.color.rgb = color
    
    p_m_desc = tf_diag.add_paragraph()
    p_m_desc.text = f"   {desc}"
    p_m_desc.font.name = 'Century Gothic'
    p_m_desc.font.size = Pt(10)
    p_m_desc.font.color.rgb = SLATE_GRAY
    p_m_desc.space_after = Pt(10)

p_eq = tf_diag.add_paragraph()
p_eq.text = "➔ DELIVERING FASTER, SAFER, CHEAPER, GREENER TRAVEL"
p_eq.font.name = 'Century Gothic'
p_eq.font.size = Pt(11)
p_eq.font.bold = True
p_eq.font.color.rgb = NEON_CYAN
p_eq.space_before = Pt(8)


# ----------------- SLIDE 4: INTELLIGENT JOURNEY PLANNING -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Key Feature 1: Multi-Objective Routing", "Intelligent Journey Planning")

# Left Column text
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.2), Inches(4.8))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_h = tf_left.paragraphs[0]
p_h.text = "Customized Dijkstra Routing Weight Calculations"
p_h.font.name = 'Century Gothic'
p_h.font.size = Pt(20)
p_h.font.bold = True
p_h.font.color.rgb = WHITE
p_h.space_after = Pt(14)

bullets = [
    ("4 Travel Objectives", "Toggle dynamically between Fastest, Cheapest, Safest, or Eco-friendly. The router re-weights graph edges based on selected objective metrics."),
    ("Multi-Modal Integration", "Connects walking, cycling, auto-rickshaw, city buses, and local metros into a unified route recommendation."),
    ("Balanced Mode Selection", "A hybrid algorithm solving for multi-objective Pareto optimization, recommending the best compromise route.")
]

for title, text in bullets:
    p_b_title = tf_left.add_paragraph()
    p_b_title.text = f"▪ {title}"
    p_b_title.font.name = 'Century Gothic'
    p_b_title.font.size = Pt(12)
    p_b_title.font.bold = True
    p_b_title.font.color.rgb = SOFT_TEAL
    
    p_b_text = tf_left.add_paragraph()
    p_b_text.text = f"  {text}"
    p_b_text.font.name = 'Century Gothic'
    p_b_text.font.size = Pt(11)
    p_b_text.font.color.rgb = SLATE_GRAY
    p_b_text.space_after = Pt(8)

# Right Column Graphic: Route comparison mockup
create_card(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.8), DARK_CARD, BORDER_BLUE)

right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.9), Inches(5.6), Inches(4.4))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_sim_t = tf_right.paragraphs[0]
p_sim_t.text = "ROUTE SIMULATION METRIC COMPARISON"
p_sim_t.font.name = 'Century Gothic'
p_sim_t.font.size = Pt(11)
p_sim_t.font.bold = True
p_sim_t.font.color.rgb = NEON_CYAN
p_sim_t.space_after = Pt(14)

sim_routes = [
    ("FASTEST ROUTE (Taxi / Ride-Share)", "14 Mins  |  ₹240 Fare  |  2.4 kg CO2  |  Safety Score: 7.2", CRIMSON_RED),
    ("SAFEST ROUTE (Metro + Walk Corridor)", "22 Mins  |  ₹40 Fare  |  0.3 kg CO2  |  Safety Score: 9.8", NEON_CYAN),
    ("ECO-FRIENDLY ROUTE (Bicycle Lane)", "30 Mins  |  ₹0 Fare  |  0.0 kg CO2  |  Safety Score: 8.5", EMERALD_GREEN)
]

for title, stats, color in sim_routes:
    p_r_title = tf_right.add_paragraph()
    p_r_title.text = f"➔ {title}"
    p_r_title.font.name = 'Century Gothic'
    p_r_title.font.size = Pt(12)
    p_r_title.font.bold = True
    p_r_title.font.color.rgb = color
    
    p_r_stats = tf_right.add_paragraph()
    p_r_stats.text = f"   {stats}"
    p_r_stats.font.name = 'Century Gothic'
    p_r_stats.font.size = Pt(11)
    p_r_stats.font.color.rgb = WHITE
    p_r_stats.space_after = Pt(12)

p_r_note = tf_right.add_paragraph()
p_r_note.text = "*Map interface powered by OpenStreetMap coordinates, rendering safety and lighting gradients."
p_r_note.font.name = 'Century Gothic'
p_r_note.font.size = Pt(9)
p_r_note.font.italic = True
p_r_note.font.color.rgb = SLATE_GRAY


# ----------------- SLIDE 5: ROUTE SAFETY & WOMEN'S SAFETY -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Key Feature 2: Route Safety & Emergency Console", "Women's Safety & Security")

# Left Column text
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.2), Inches(4.8))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_h = tf_left.paragraphs[0]
p_h.text = "Guardian AI: Continuous Safety Assessment"
p_h.font.name = 'Century Gothic'
p_h.font.size = Pt(20)
p_h.font.bold = True
p_h.font.color.rgb = WHITE
p_h.space_after = Pt(14)

safety_features = [
    ("AI Safety Score (0-10)", "Calculates safety ratings based on road lighting coverage, reported incidents, historical crime coordinates, and community feedback."),
    ("Women's Security Module", "One-Tap Panic Alarm, live coordinates sharing with trusted contacts, SHE Teams quick access integration, and simulated audio call deterrents."),
    ("Emergency Locator", "Instant query maps highlighting the nearest active police stations, municipal hospitals, and quick security response units.")
]

for title, text in safety_features:
    p_f_title = tf_left.add_paragraph()
    p_f_title.text = f"▪ {title}"
    p_f_title.font.name = 'Century Gothic'
    p_f_title.font.size = Pt(12)
    p_f_title.font.bold = True
    p_f_title.font.color.rgb = CRIMSON_RED
    
    p_f_text = tf_left.add_paragraph()
    p_f_text.text = f"  {text}"
    p_f_text.font.name = 'Century Gothic'
    p_f_text.font.size = Pt(11)
    p_f_text.font.color.rgb = SLATE_GRAY
    p_f_text.space_after = Pt(8)

# Right Column Graphic: SOS console mockup
create_card(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.8), DARK_CARD, CRIMSON_RED)

right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.9), Inches(5.6), Inches(4.4))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_sos_t = tf_right.paragraphs[0]
p_sos_t.text = "GUARDIAN SOS EMERGENCY SYSTEM"
p_sos_t.font.name = 'Century Gothic'
p_sos_t.font.size = Pt(11)
p_sos_t.font.bold = True
p_sos_t.font.color.rgb = WHITE
p_sos_t.space_after = Pt(14)

sos_console = [
    ("ONE-TAP PANIC BUTTON", "Sends encrypted SMS containing live GPS coordinates to trusted contacts and municipal emergency response cells.", CRIMSON_RED),
    ("VOICE-TRIGGERED ALARM", "Monitors device microphone for specific distress keywords. Triggers automated emergency broadcast silently if activated.", WARNING_GOLD),
    ("LIVE PATH MONITORING", "Compares active route telemetry with predicted path; triggers warnings if the user deviates into low-safety zones.", SOFT_TEAL)
]

for title, desc, color in sos_console:
    p_c_title = tf_right.add_paragraph()
    p_c_title.text = f"🚨 {title}"
    p_c_title.font.name = 'Century Gothic'
    p_c_title.font.size = Pt(12)
    p_c_title.font.bold = True
    p_c_title.font.color.rgb = color
    
    p_c_desc = tf_right.add_paragraph()
    p_c_desc.text = f"   {desc}"
    p_c_desc.font.name = 'Century Gothic'
    p_c_desc.font.size = Pt(10.5)
    p_c_desc.font.color.rgb = SLATE_GRAY
    p_c_desc.space_after = Pt(10)


# ----------------- SLIDE 6: AI TRAFFIC & CARBON TRACKING -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Key Feature 3: Congestion Prediction & Carbon Tracker", "AI Traffic & Sustainability")

# Left Column text
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.2), Inches(4.8))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_h = tf_left.paragraphs[0]
p_h.text = "Preemptive Traffic Analytics & Carbon Accounting"
p_h.font.name = 'Century Gothic'
p_h.font.size = Pt(20)
p_h.font.bold = True
p_h.font.color.rgb = WHITE
p_h.space_after = Pt(14)

traffic_features = [
    ("Preemptive Traffic Engine", "Uses a machine learning Decision Tree regressor to forecast congestion density at nodes 30–60 minutes in advance, suggesting proactive rerouting."),
    ("Carbon Offset Calculator", "Compares CO2 footprint against driving a single-occupancy vehicle. Shows exact emissions saved by choosing walking, cycling, or metro."),
    ("Green Commute XP Rewards", "Commuters earn points (XP) for sustainable routes, unlocking gamified badges and redeemable city transit vouchers.")
]

for title, text in traffic_features:
    p_t_title = tf_left.add_paragraph()
    p_t_title.text = f"▪ {title}"
    p_t_title.font.name = 'Century Gothic'
    p_t_title.font.size = Pt(12)
    p_t_title.font.bold = True
    p_t_title.font.color.rgb = EMERALD_GREEN
    
    p_t_text = tf_left.add_paragraph()
    p_t_text.text = f"  {text}"
    p_t_text.font.name = 'Century Gothic'
    p_t_text.font.size = Pt(11)
    p_t_text.font.color.rgb = SLATE_GRAY
    p_t_text.space_after = Pt(8)

# Right Column Graphic
create_card(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.8), DARK_CARD, EMERALD_GREEN)

right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.9), Inches(5.6), Inches(4.4))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_eco_t = tf_right.paragraphs[0]
p_eco_t.text = "SUSTAINABILITY & TELEMETRY INSIGHTS"
p_eco_t.font.name = 'Century Gothic'
p_eco_t.font.size = Pt(11)
p_eco_t.font.bold = True
p_eco_t.font.color.rgb = WHITE
p_eco_t.space_after = Pt(14)

p_offset_t = tf_right.add_paragraph()
p_offset_t.text = "🌲 User Carbon Offset Ledger"
p_offset_t.font.name = 'Century Gothic'
p_offset_t.font.size = Pt(13)
p_offset_t.font.bold = True
p_offset_t.font.color.rgb = EMERALD_GREEN

p_offset_d = tf_right.add_paragraph()
p_offset_d.text = "   Total Savings: 82.4 kg CO2 Saved (Equivalent to 4.1 saplings planted)\n   Commute XP Balance: 1,240 XP  |  Sustain Grade: A"
p_offset_d.font.name = 'Century Gothic'
p_offset_d.font.size = Pt(11)
p_offset_d.font.color.rgb = WHITE
p_offset_d.space_after = Pt(14)

p_alert_t = tf_right.add_paragraph()
p_alert_t.text = "⚠️ Preemptive Congestion Warning"
p_alert_t.font.name = 'Century Gothic'
p_alert_t.font.size = Pt(13)
p_alert_t.font.bold = True
p_alert_t.font.color.rgb = WARNING_GOLD

p_alert_d = tf_right.add_paragraph()
p_alert_d.text = "   Junction 4 (Main Road Hub): Gridlock predicted at 18:00.\n   Recommendation: Divert to Flyover Bypass to save 12 mins and reduce idling emissions by 0.4 kg CO2."
p_alert_d.font.name = 'Century Gothic'
p_alert_d.font.size = Pt(11)
p_alert_d.font.color.rgb = WHITE


# ----------------- SLIDE 7: SMART COST OPTIMIZATION -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Key Feature 4: Cost vs. Time Analysis", "Smart Cost Optimizer")

# Left Column text
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.2), Inches(4.8))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_h = tf_left.paragraphs[0]
p_h.text = "Tackling Transportation Inflation"
p_h.font.name = 'Century Gothic'
p_h.font.size = Pt(20)
p_h.font.bold = True
p_h.font.color.rgb = WHITE
p_h.space_after = Pt(14)

cost_features = [
    ("Fare Comparison Engine", "Aggregates ticketing fares across public metros, local city buses, auto-rickshaws, and ride-hailing services in real-time."),
    ("Cost vs. Time Trade-off Graph", "Shows the Pareto Frontier of commute choices, identifying where a small increase in travel time yields a high cost saving."),
    ("Alternative Mode Prompts", "Triggers suggestions to swap modes (e.g., walking 300 meters to connect with a metro node rather than ordering a high-surge cab).")
]

for title, text in cost_features:
    p_c_title = tf_left.add_paragraph()
    p_c_title.text = f"▪ {title}"
    p_c_title.font.name = 'Century Gothic'
    p_c_title.font.size = Pt(12)
    p_c_title.font.bold = True
    p_c_title.font.color.rgb = SOFT_TEAL
    
    p_c_text = tf_left.add_paragraph()
    p_c_text.text = f"  {text}"
    p_c_text.font.name = 'Century Gothic'
    p_c_text.font.size = Pt(11)
    p_c_text.font.color.rgb = SLATE_GRAY
    p_c_text.space_after = Pt(8)

# Right Column Graphic
create_card(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.8), DARK_CARD, BORDER_BLUE)

right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.9), Inches(5.6), Inches(4.4))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_cost_t = tf_right.paragraphs[0]
p_cost_t.text = "FARE OPTIMIZATION ENGINE METRICS"
p_cost_t.font.name = 'Century Gothic'
p_cost_t.font.size = Pt(11)
p_cost_t.font.bold = True
p_cost_t.font.color.rgb = NEON_CYAN
p_cost_t.space_after = Pt(14)

p_save_t = tf_right.add_paragraph()
p_save_t.text = "💰 Projected Monthly Savings"
p_save_t.font.name = 'Century Gothic'
p_save_t.font.size = Pt(13)
p_save_t.font.bold = True
p_save_t.font.color.rgb = NEON_CYAN

p_save_d = tf_right.add_paragraph()
p_save_d.text = "   User Average: ₹2,840 Saved/Month\n   Achieved by: Swapping 35% of peak taxi rides to metro connections."
p_save_d.font.name = 'Century Gothic'
p_save_d.font.size = Pt(11)
p_save_d.font.color.rgb = WHITE
p_save_d.space_after = Pt(14)

p_prompt_t = tf_right.add_paragraph()
p_prompt_t.text = "💡 Dynamic Savings Suggestion"
p_prompt_t.font.name = 'Century Gothic'
p_prompt_t.font.size = Pt(13)
p_prompt_t.font.bold = True
p_prompt_t.font.color.rgb = WHITE

p_prompt_d = tf_right.add_paragraph()
p_prompt_d.text = "   Cab fare is currently surging (₹240, 14 mins).\n   Action: Walk 5 mins to Metro Station 2, take Metro Line 3.\n   Result: Fare drops to ₹40 (Savings: ₹200). Travel time increases by only 8 mins."
p_prompt_d.font.name = 'Century Gothic'
p_prompt_d.font.size = Pt(11)
p_prompt_d.font.color.rgb = SLATE_GRAY


# ----------------- SLIDE 8: SYSTEM ARCHITECTURE & TECH STACK -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Robust Technology Stack & Layered Design", "System Engineering")

# 5 Columns representing tech stack layers
tech_cols = [
    {
        "title": "🖥️ CLIENT DEV",
        "label": "Frontend",
        "tech": "Vite\nReact\nTailwind CSS\nLucide Icons\nFramer Motion",
        "color": SOFT_TEAL
    },
    {
        "title": "⚙️ ENGINE API",
        "label": "Backend",
        "tech": "FastAPI\nUvicorn Server\nJWT Security\nWebSockets for SOS broadcasts",
        "color": SOFT_TEAL
    },
    {
        "title": "🗄️ SPATIAL DB",
        "label": "Database",
        "tech": "PostgreSQL\nPostGIS Ext\nSQLAlchemy ORM\nIncident Seeding",
        "color": SOFT_TEAL
    },
    {
        "title": "🗺️ MAP ENGINE",
        "label": "GIS Routing",
        "tech": "OpenStreetMap\nOSRM Route Solver\nLeaflet Map APIs\nGeoJSON parsing",
        "color": SOFT_TEAL
    },
    {
        "title": "🧠 AI MODELS",
        "label": "ML Inference",
        "tech": "Scikit-Learn\nDecision Trees\nDijkstra Graph Solver\nNLP Assistant",
        "color": SOFT_TEAL
    }
]

col_w = Inches(2.1)
gap_w = Inches(0.3)
left_m = Inches(0.8)
top_m = Inches(1.8)
col_h = Inches(4.3)

for idx, tech in enumerate(tech_cols):
    col_left = left_m + idx * (col_w + gap_w)
    create_card(slide, col_left, top_m, col_w, col_h, DARK_CARD, BORDER_BLUE)
    
    box = slide.shapes.add_textbox(col_left + Inches(0.1), top_m + Inches(0.15), col_w - Inches(0.2), col_h - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.text = tech["title"]
    p_t.font.name = 'Century Gothic'
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    
    p_l = tf.add_paragraph()
    p_l.text = tech["label"].upper()
    p_l.font.name = 'Century Gothic'
    p_l.font.size = Pt(9)
    p_l.font.bold = True
    p_l.font.color.rgb = NEON_CYAN
    p_l.space_after = Pt(14)
    
    p_lst = tf.add_paragraph()
    p_lst.text = tech["tech"]
    p_lst.font.name = 'Century Gothic'
    p_lst.font.size = Pt(11)
    p_lst.font.color.rgb = SLATE_GRAY
    
    # Adjust line spacing for tech list
    p_lst.line_spacing = 1.3

# Bottom architecture explanation note
note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5))
tf_note = note_box.text_frame
p_note = tf_note.paragraphs[0]
p_note.text = "▪ System integration: The Graph router fetches active nodes from OpenRouteService. Dijkstra edge cost calculation uses ML traffic speed predictions and lighting weights."
p_note.font.name = 'Century Gothic'
p_note.font.size = Pt(10.5)
p_note.font.color.rgb = SLATE_GRAY


# ----------------- SLIDE 9: URBAN MOBILITY ANALYTICS DASHBOARD -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Urban Mobility Analytics Dashboard", "Smart City Planning")

# Left Column text
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.2), Inches(4.8))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_h = tf_left.paragraphs[0]
p_h.text = "Aggregating Telemetry for Municipal Planners"
p_h.font.name = 'Century Gothic'
p_h.font.size = Pt(20)
p_h.font.bold = True
p_h.font.color.rgb = WHITE
p_h.space_after = Pt(14)

dash_points = [
    ("Congestion Heatmapping", "Helps municipal authorities identify bottlenecks and peak rush-hour patterns to optimize public transit scheduling."),
    ("Incident Density Reports", "Aggregates crowd-sourced citizen hazard logs (unlit roads, potholes) to direct city maintenance crews."),
    ("Environmental Impact Tracker", "Calculates cumulative carbon emissions saved city-wide, supporting sustainability metrics and green policies.")
]

for title, text in dash_points:
    p_d_title = tf_left.add_paragraph()
    p_d_title.text = f"▪ {title}"
    p_d_title.font.name = 'Century Gothic'
    p_d_title.font.size = Pt(12)
    p_d_title.font.bold = True
    p_d_title.font.color.rgb = SOFT_TEAL
    
    p_d_text = tf_left.add_paragraph()
    p_d_text.text = f"  {text}"
    p_d_text.font.name = 'Century Gothic'
    p_d_text.font.size = Pt(11)
    p_d_text.font.color.rgb = SLATE_GRAY
    p_d_text.space_after = Pt(8)

# Right Column Graphic: Mock dashboard charts
create_card(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(4.8), DARK_CARD, BORDER_BLUE)

right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.9), Inches(5.6), Inches(4.4))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_dash_t = tf_right.paragraphs[0]
p_dash_t.text = "ADMINISTRATOR CENTRAL DASHBOARD"
p_dash_t.font.name = 'Century Gothic'
p_dash_t.font.size = Pt(11)
p_dash_t.font.bold = True
p_dash_t.font.color.rgb = NEON_CYAN
p_dash_t.space_after = Pt(14)

# Admin indicators mockup text
p_stats_h = tf_right.add_paragraph()
p_stats_h.text = "📊 Today's Cumulative Metrics (City-wide)"
p_stats_h.font.name = 'Century Gothic'
p_stats_h.font.size = Pt(13)
p_stats_h.font.bold = True
p_stats_h.font.color.rgb = WHITE
p_stats_h.space_after = Pt(8)

p_stats_d = tf_right.add_paragraph()
p_stats_d.text = "   ▪ Active Commuters: 1,240 Citizens Online\n   ▪ Verified Safety Incidents Logged: 4 Resolved\n   ▪ CO2 Footprint Reduction: 142.8 kg Saved Today"
p_stats_d.font.name = 'Century Gothic'
p_stats_d.font.size = Pt(11.5)
p_stats_d.font.color.rgb = SLATE_GRAY
p_stats_d.space_after = Pt(14)

p_modes_h = tf_right.add_paragraph()
p_modes_h.text = "🚲 Travel Mode Distribution"
p_modes_h.font.name = 'Century Gothic'
p_modes_h.font.size = Pt(13)
p_modes_h.font.bold = True
p_modes_h.font.color.rgb = WHITE
p_modes_h.space_after = Pt(8)

p_modes_d = tf_right.add_paragraph()
p_modes_d.text = "   ▪ Metro / City Bus: 35% usage share\n   ▪ Bicycling / Walking corridors: 45% usage share\n   ▪ Single-Occupancy Taxi / Auto: 20% usage share"
p_modes_d.font.name = 'Century Gothic'
p_modes_d.font.size = Pt(11.5)
p_modes_d.font.color.rgb = SLATE_GRAY


# ----------------- SLIDE 10: EXPECTED IMPACT & FUTURE VISION -----------------
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_header(slide, "Empowering the Smart Cities of Tomorrow", "Expected Impact & Vision")

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
tf_sub = sub_box.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "A sustainable, data-driven approach to urban commuting."
p_sub.font.name = 'Century Gothic'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = SLATE_GRAY

# 4 KPI cards side by side (impact numbers)
kpis = [
    {"num": "100%", "label": "UNIFIED PLATFORM", "desc": "Bridges maps, emergency SOS, and carbon tracker apps into a single view.", "color": SOFT_TEAL},
    {"num": "-40%", "label": "SAFETY INCIDENTS", "desc": "Reduces night travel risks via light-weighted routing algorithms.", "color": CRIMSON_RED},
    {"num": "-25%", "label": "CARBON EMISSIONS", "desc": "Promotes green transit through gamified XP reward frameworks.", "color": EMERALD_GREEN},
    {"num": "3.5x", "label": "TRANSIT ADOPTION", "desc": "Boosts public transit utilization via integrated cost-savings analysis.", "color": WARNING_GOLD}
]

col_w = Inches(2.7)
gap = Inches(0.3)
left_margin = Inches(0.8)
top_margin = Inches(2.1)
col_h = Inches(2.2)

for idx, kpi in enumerate(kpis):
    col_left = left_margin + idx * (col_w + gap)
    create_card(slide, col_left, top_margin, col_w, col_h, DARK_CARD, BORDER_BLUE)
    
    box = slide.shapes.add_textbox(col_left + Inches(0.1), top_margin + Inches(0.1), col_w - Inches(0.2), col_h - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    
    p_num = tf.paragraphs[0]
    p_num.text = kpi["num"]
    p_num.font.name = 'Century Gothic'
    p_num.font.size = Pt(28)
    p_num.font.bold = True
    p_num.font.color.rgb = kpi["color"]
    p_num.alignment = PP_ALIGN.CENTER
    
    p_lbl = tf.add_paragraph()
    p_lbl.text = kpi["label"]
    p_lbl.font.name = 'Century Gothic'
    p_lbl.font.size = Pt(10)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = WHITE
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.space_after = Pt(8)
    
    p_desc = tf.add_paragraph()
    p_desc.text = kpi["desc"]
    p_desc.font.name = 'Century Gothic'
    p_desc.font.size = Pt(10)
    p_desc.font.color.rgb = SLATE_GRAY
    p_desc.alignment = PP_ALIGN.CENTER

# Future Milestones Card (takes up the bottom half)
create_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.8), DARK_CARD, BORDER_BLUE)

future_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.6))
tf_future = future_box.text_frame
tf_future.word_wrap = True

p_f_h = tf_future.paragraphs[0]
p_f_h.text = "THE ROAD AHEAD: FUTURE MILESTONES"
p_f_h.font.name = 'Century Gothic'
p_f_h.font.size = Pt(12)
p_f_h.font.bold = True
p_f_h.font.color.rgb = NEON_CYAN
p_f_h.space_after = Pt(8)

p_f_d = tf_future.add_paragraph()
p_f_d.text = "▪ Phase 2 IoT Integration: Connect municipal smart streetlights to route safety indices dynamically.\n▪ Smart Grid Traffic Signals: Implement real-time predictive green signal tuning based on congestion density forecasts.\n▪ SHE Teams Live Integration: Form direct secure telemetry links with active SHE Teams patrol vehicles for automated panic intercepts.\n▪ Multi-City Expansion: Configure routing models to adapt dynamically to regional transit data API feeds."
p_f_d.font.name = 'Century Gothic'
p_f_d.font.size = Pt(11)
p_f_d.font.color.rgb = SLATE_GRAY

# Save presentation
output_path = os.path.join(os.path.dirname(__file__), "UrbanFlow_AI_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}!")
